"""Briefing export to PPTX format."""

import base64
import io
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from foothold_sitac.briefing import Briefing, Homeplate, MissionType, Objective, Package


def _rgb(r: int, g: int, b: int) -> Any:
    """Create RGBColor with proper typing."""
    return RGBColor(r, g, b)  # type: ignore[no-untyped-call]


# Color scheme (dark theme matching the app)
COLOR_BG_DARK: Any = _rgb(0x1A, 0x1A, 0x2E)
COLOR_TEXT_LIGHT: Any = _rgb(0xE0, 0xE0, 0xE0)
COLOR_ACCENT: Any = _rgb(0x4A, 0x90, 0xD9)
COLOR_HEADER: Any = _rgb(0x2D, 0x2D, 0x44)

# Mission type colors (unused for now but available for future enhancements)
MISSION_COLORS: dict[MissionType, Any] = {
    MissionType.CAP: _rgb(0x4A, 0x90, 0xD9),  # Blue
    MissionType.SEAD: _rgb(0xF5, 0xA6, 0x23),  # Orange
    MissionType.DEAD: _rgb(0xE5, 0x4B, 0x4B),  # Red
    MissionType.CAS: _rgb(0x4E, 0xCB, 0x71),  # Green
    MissionType.STRIKE: _rgb(0x9B, 0x59, 0xB6),  # Purple
    MissionType.SWEEP: _rgb(0x3B, 0xB0, 0xD9),  # Light blue
    MissionType.ESCORT: _rgb(0xF3, 0x9C, 0x12),  # Gold
    MissionType.RECCE: _rgb(0x95, 0xA5, 0xA6),  # Gray
    MissionType.TRANSPORT: _rgb(0x8B, 0x45, 0x13),  # Brown
    MissionType.CSAR: _rgb(0xFF, 0xD7, 0x00),  # Yellow/Gold
}


def create_briefing_pptx(briefing: Briefing, map_image_base64: str | None = None) -> io.BytesIO:
    """Generate a PPTX presentation from a briefing.

    Args:
        briefing: The briefing data to export.
        map_image_base64: Optional base64-encoded PNG image of the tactical map.

    Returns:
        BytesIO buffer containing the PPTX file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Add slides
    _add_title_slide(prs, briefing)

    if map_image_base64:
        _add_map_slide(prs, map_image_base64)

    if briefing.situation:
        _add_text_slide(prs, "Situation", briefing.situation)

    if briefing.homeplates:
        _add_airbases_slide(prs, briefing.homeplates)

    if briefing.objectives:
        _add_objectives_slide(prs, briefing.objectives)

    for package in briefing.packages:
        _add_package_slide(prs, package, briefing.homeplates)

    if briefing.weather:
        _add_text_slide(prs, "Weather", briefing.weather)

    if briefing.comms_plan:
        _add_text_slide(prs, "Comms Plan", briefing.comms_plan)

    if briefing.notes:
        _add_text_slide(prs, "Notes", briefing.notes)

    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer


def _set_slide_background(slide: Any, color: Any) -> None:
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_title_slide(prs: Any, briefing: Briefing) -> None:
    """Add title slide with briefing metadata."""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = briefing.title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_LIGHT
    p.alignment = PP_ALIGN.CENTER

    # Server badge
    server_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.5))
    tf = server_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Server: {briefing.server_name}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # Date/time
    timing_parts = []
    if briefing.mission_date:
        timing_parts.append(f"Date: {briefing.mission_date}")
    if briefing.mission_time:
        timing_parts.append(f"Time: {briefing.mission_time}")

    if timing_parts:
        timing_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(12.333), Inches(0.5))
        tf = timing_box.text_frame
        p = tf.paragraphs[0]
        p.text = "  |  ".join(timing_parts)
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER


def _add_map_slide(prs: Any, map_image_base64: str) -> None:
    """Add full-slide tactical map image."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Decode image
    image_data = base64.b64decode(map_image_base64)
    image_stream = io.BytesIO(image_data)

    # Add image centered and scaled to fit
    slide.shapes.add_picture(
        image_stream,
        Inches(0),
        Inches(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )


def _add_text_slide(prs: Any, title: str, content: str) -> None:
    """Add a simple text content slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(18)
    p.font.color.rgb = COLOR_TEXT_LIGHT


def _add_airbases_slide(prs: Any, homeplates: list[Homeplate]) -> None:
    """Add airbases table slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Airbases"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # Table
    rows = len(homeplates) + 1  # +1 for header
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(0.5 * rows)).table

    # Set column widths
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(3.833)

    # Header row
    headers = ["Name", "TACAN", "Runway", "Frequencies"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_HEADER
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, hp in enumerate(homeplates, start=1):
        row_data = [
            hp.name,
            hp.tacan or "-",
            f"{hp.runway_heading}°" if hp.runway_heading else "-",
            ", ".join(hp.frequencies) if hp.frequencies else "-",
        ]
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_BG_DARK
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_TEXT_LIGHT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_objectives_slide(prs: Any, objectives: list[Objective]) -> None:
    """Add objectives slide with mission badges."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Objectives"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # Sort by priority
    sorted_objectives = sorted(objectives, key=lambda o: o.priority)

    y_pos = 1.5
    for obj in sorted_objectives:
        # Priority badge and zone name
        obj_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(0.5))
        tf = obj_box.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"P{obj.priority}"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = _rgb(0xF5, 0xA6, 0x23)  # Orange for priority

        run = p.add_run()
        run.text = f"  {obj.zone_name}"
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_TEXT_LIGHT

        # Mission types
        if obj.mission_requirements:
            mission_text = "   [" + " | ".join(mt.value for mt in obj.mission_requirements) + "]"
            run = p.add_run()
            run.text = mission_text
            run.font.size = Pt(14)
            run.font.color.rgb = COLOR_ACCENT

        y_pos += 0.4

        # Notes
        if obj.notes:
            notes_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos), Inches(11.833), Inches(0.4))
            tf = notes_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = obj.notes
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = _rgb(0xA0, 0xA0, 0xA0)
            y_pos += 0.4

        y_pos += 0.2


def _add_package_slide(prs: Any, package: Package, homeplates: list[Homeplate]) -> None:
    """Add package slide with flights table."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_background(slide, COLOR_BG_DARK)

    # Package name
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = package.name
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # Target zone and mission type
    subtitle_parts = []
    if package.target_zone:
        subtitle_parts.append(f"Target: {package.target_zone}")
    if package.mission_type:
        subtitle_parts.append(f"Mission: {package.mission_type.value}")

    if subtitle_parts:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(10), Inches(0.4))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = " | ".join(subtitle_parts)
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT_LIGHT

    # Build homeplate lookup
    hp_lookup = {hp.id: hp.name for hp in homeplates}

    # Flights table
    if package.flights:
        rows = len(package.flights) + 1
        cols = 8
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.4 * rows)).table

        # Set column widths
        table.columns[0].width = Inches(2)  # Callsign
        table.columns[1].width = Inches(1.5)  # Aircraft
        table.columns[2].width = Inches(0.5)  # #
        table.columns[3].width = Inches(1.2)  # Mission
        table.columns[4].width = Inches(2)  # Departure
        table.columns[5].width = Inches(2)  # Arrival
        table.columns[6].width = Inches(1.1)  # PUSH
        table.columns[7].width = Inches(1.1)  # TOT

        # Header
        headers = ["Callsign", "Aircraft", "#", "Mission", "Departure", "Arrival", "PUSH", "TOT"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_HEADER
            p = cell.text_frame.paragraphs[0]
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_TEXT_LIGHT
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for row_idx, flight in enumerate(package.flights, start=1):
            departure_name = hp_lookup.get(flight.departure_id, "-") if flight.departure_id else "-"
            arrival_name = hp_lookup.get(flight.arrival_id, departure_name) if flight.arrival_id else departure_name

            row_data = [
                flight.callsign,
                flight.aircraft_type,
                str(flight.num_aircraft),
                flight.mission_type.value,
                departure_name,
                arrival_name,
                flight.push_time or "-",
                flight.tot or "-",
            ]
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BG_DARK
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_TEXT_LIGHT
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Package notes
    if package.notes:
        y_pos = 1.8 + (0.4 * (len(package.flights) + 1)) + 0.3 if package.flights else 1.8
        notes_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12.333), Inches(1))
        tf = notes_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = package.notes
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = _rgb(0xA0, 0xA0, 0xA0)
